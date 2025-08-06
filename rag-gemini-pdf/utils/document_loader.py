import fitz  # PyMuPDF-used to read and extract text from PDF
import docx #from python-docx library used to read DOCX files 
from pptx import Presentation # From python-pptx library; used to load PowerPoint files and access slide contents.
import os #Standard Python module for handling file paths and extensions.

def load_unstructured_file(file_path):
    #Splits the file path into (filename, extension)
    ext = os.path.splitext(file_path)[1].lower() #[1].lower(): Extracts the extension (e.g., .pdf) and converts it to lowercase for consistency.
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else: #if not a supported file structure, returns empty string
        return ""

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path) # Opens the PDF using PyMuPDF.
    text = "\n".join(page.get_text() for page in doc)
    #page.get_text(): Extracts all text from each page.
    #"\n".join(...): Joins text from all pages with newlines to 
    # create a single string of the full document content.
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path) #Opens the Word document.
    return "\n".join([p.text for p in doc.paragraphs])#Loops through each paragraph and extracts the text.
    #"\n".join(...): Combines all paragraph texts into a single string with line breaks.


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    texts = []#Initializes an empty list to collect texts from all slides.
    for slide in prs.slides:
        for shape in slide.shapes:#Iterates through each slide and each shape (like textbox, title, etc.).
            if hasattr(shape, "text"):#Some shapes may not have text; this safely checks if text exists.
                texts.append(shape.text)#Collects text from valid shapes.
    return "\n".join(texts) #Combines all extracted slide texts into one big string separated by line breaks.


# Sentence-based Fixed-size Chunking
def chunk_text(text, max_tokens=300):
    sentences = text.split(". ") #Splits the full text into sentences using . as the separator.
    chunks, current_chunk = [], ""
    for sentence in sentences:
        if len((current_chunk + sentence).split()) <= max_tokens:
            current_chunk += sentence + ". "
        else:
            chunks.append(current_chunk.strip())
            current_chunk = sentence + ". "
    if current_chunk:
        chunks.append(current_chunk.strip())#.strip() removes leading/trailing spaces.
    return chunks
